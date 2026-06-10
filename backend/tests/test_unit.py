"""
MIND FORGE — Backend Unit Tests
Tests pure business logic: MPIN hashing, JWT, and input validation schemas.
No database, no network, no Redis required.
Run: cd backend && python3 -m pytest tests/test_unit.py -v
"""

import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

from datetime import timedelta, datetime, timezone

import pytest
from jose import JWTError

# ── MPIN hashing ──────────────────────────────────────────────────────────────

from app.core.security import hash_mpin, verify_mpin


def test_hash_mpin_returns_string():
    h = hash_mpin("123456")
    assert isinstance(h, str)
    assert len(h) > 0


def test_verify_mpin_correct():
    h = hash_mpin("123456")
    assert verify_mpin("123456", h) is True


def test_verify_mpin_wrong():
    h = hash_mpin("123456")
    assert verify_mpin("654321", h) is False


def test_hash_mpin_different_each_time():
    # bcrypt uses a random salt — two hashes of the same mpin must differ
    h1 = hash_mpin("000000")
    h2 = hash_mpin("000000")
    assert h1 != h2
    # But both must still verify
    assert verify_mpin("000000", h1)
    assert verify_mpin("000000", h2)


# ── JWT encode / decode ───────────────────────────────────────────────────────

from app.core.security import create_access_token, create_refresh_token, decode_access_token


def test_jwt_roundtrip_preserves_claims():
    token = create_access_token({"sub": "42", "role": "teacher"})
    payload = decode_access_token(token)
    assert payload["sub"] == "42"
    assert payload["role"] == "teacher"
    assert payload["type"] == "access"


def test_refresh_token_has_correct_type():
    token = create_refresh_token({"sub": "7", "role": "student"})
    payload = decode_access_token(token)
    assert payload["type"] == "refresh"
    assert payload["sub"] == "7"


def test_expired_access_token_raises():
    token = create_access_token(
        {"sub": "1"}, expires_delta=timedelta(seconds=-1)
    )
    with pytest.raises(JWTError):
        decode_access_token(token)


def test_tampered_token_raises():
    token = create_access_token({"sub": "1", "role": "teacher"})
    # Flip the last character to invalidate the signature
    tampered = token[:-1] + ("A" if token[-1] != "A" else "B")
    with pytest.raises(JWTError):
        decode_access_token(tampered)


def test_fake_token_raises():
    with pytest.raises(JWTError):
        decode_access_token("not.a.real.token")


def test_access_token_expiry_field_is_in_future():
    token = create_access_token({"sub": "1"})
    payload = decode_access_token(token)
    exp = datetime.fromtimestamp(payload["exp"], tz=timezone.utc)
    assert exp > datetime.now(timezone.utc)


# ── Input validation (Pydantic schemas) ──────────────────────────────────────

from pydantic import ValidationError
from app.schemas.user import UserLoginRequest, UserRegisterRequest
from app.models.user import UserRole


class TestUserLoginRequest:
    def test_valid_request_passes(self):
        req = UserLoginRequest(username="alice", mpin="123456")
        assert req.username == "alice"
        assert req.mpin == "123456"

    def test_mpin_must_be_6_digits(self):
        with pytest.raises(ValidationError):
            UserLoginRequest(username="alice", mpin="12345")  # 5 digits

    def test_mpin_must_be_numeric(self):
        with pytest.raises(ValidationError):
            UserLoginRequest(username="alice", mpin="12345a")

    def test_mpin_too_long_fails(self):
        with pytest.raises(ValidationError):
            UserLoginRequest(username="alice", mpin="1234567")

    def test_null_byte_in_username_rejected(self):
        with pytest.raises(ValidationError):
            UserLoginRequest(username="admin\x00injected", mpin="123456")

    def test_oversized_username_rejected(self):
        with pytest.raises(ValidationError):
            UserLoginRequest(username="A" * 151, mpin="123456")

    def test_username_at_max_length_accepted(self):
        req = UserLoginRequest(username="A" * 150, mpin="123456")
        assert len(req.username) == 150

    def test_missing_mpin_raises(self):
        with pytest.raises(ValidationError):
            UserLoginRequest(username="alice")  # type: ignore[call-arg]

    def test_missing_username_raises(self):
        with pytest.raises(ValidationError):
            UserLoginRequest(mpin="123456")  # type: ignore[call-arg]


class TestUserRegisterRequest:
    def _valid(self, **kwargs):
        # "123456" intentionally avoided here — UserRegisterRequest.mpin uses
        # the strong-MPIN validator which rejects trivial sequences. "482917"
        # has random-looking digits with no all-same / sequential / repeating
        # structure.
        defaults = dict(username="bob", mpin="482917", role=UserRole.student)
        defaults.update(kwargs)
        return UserRegisterRequest(**defaults)

    def test_valid_student_passes(self):
        req = self._valid()
        assert req.role == UserRole.student

    def test_username_too_short_fails(self):
        with pytest.raises(ValidationError):
            self._valid(username="ab")

    def test_username_too_long_fails(self):
        with pytest.raises(ValidationError):
            self._valid(username="x" * 101)

    def test_invalid_grade_rejected(self):
        with pytest.raises(ValidationError):
            self._valid(grade=7)

    def test_valid_grades_accepted(self):
        for g in (8, 9, 10):
            req = self._valid(grade=g)
            assert req.grade == g

    def test_invalid_subject_rejected(self):
        with pytest.raises(ValidationError):
            self._valid(additional_subjects=["maths"])  # not in VALID_SUBJECTS

    def test_valid_subjects_accepted(self):
        req = self._valid(additional_subjects=["economics", "computer"])
        assert "economics" in req.additional_subjects

    def test_invalid_phone_rejected(self):
        with pytest.raises(ValidationError):
            self._valid(phone="not-a-phone!!!")

    def test_valid_phone_accepted(self):
        req = self._valid(phone="+91 98765 43210")
        assert req.phone is not None

    def test_invalid_email_rejected(self):
        with pytest.raises(ValidationError):
            self._valid(email="notanemail")

    def test_valid_email_accepted(self):
        req = self._valid(email="test@example.com")
        assert req.email == "test@example.com"


# ── Uploaded-deck (.pptx) parsing ─────────────────────────────────────────────

import io as _io
from pptx import Presentation as _Presentation

from app.services.pptx_service import parse_pptx, recommended_periods


def _make_pptx(slides) -> bytes:
    """Build a minimal .pptx in memory. `slides` is a list of (title, [bullets]).
    Uses the "Title and Content" layout (index 1) so a title placeholder exists.
    """
    prs = _Presentation()
    layout = prs.slide_layouts[1]
    for title, bullets in slides:
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = title
        body = s.placeholders[1].text_frame
        body.text = bullets[0]
        for b in bullets[1:]:
            body.add_paragraph().text = b
    buf = _io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestPptxParsing:
    def test_parses_titles_and_bullets(self):
        data = _make_pptx([
            ("Photosynthesis", ["Light reaction", "Calvin cycle"]),
            ("Respiration", ["Glycolysis"]),
        ])
        slides = parse_pptx(data)
        assert len(slides) == 2
        assert slides[0]["title"] == "Photosynthesis"
        assert "- Light reaction" in slides[0]["body_md"]
        assert "- Calvin cycle" in slides[0]["body_md"]
        assert slides[1]["title"] == "Respiration"

    def test_rejects_non_pptx_bytes(self):
        with pytest.raises(ValueError):
            parse_pptx(b"this is not a pptx file")

    def test_recommended_periods_ceils_at_eight_per_period(self):
        assert recommended_periods(8, 8) == 1
        assert recommended_periods(9, 8) == 2
        assert recommended_periods(16, 8) == 2
        assert recommended_periods(17, 8) == 3

    def test_recommended_periods_never_below_one(self):
        assert recommended_periods(0, 8) == 1
        assert recommended_periods(3, 8) == 1
