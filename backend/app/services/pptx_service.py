"""
Parse an uploaded PowerPoint (.pptx) deck into MindForge slides.

Unlike presentation_service (which asks Gemini to *write* slides from a
chapter PDF), this module reads the slides the teacher already authored in
PowerPoint / Google Slides (exported as .pptx) and maps each one onto the
same {title, body_md, speaker_notes} shape the rest of the presentation
feature expects. No AI is involved — it's a fast, local, deterministic parse.

Google Slides decks are supported via File → Download → Microsoft PowerPoint
(.pptx); native .gslides files can't be uploaded directly (they live in
Drive), and legacy binary .ppt is not OpenXML so python-pptx can't read it.
"""

import io
import logging
from typing import List

from pptx import Presentation

logger = logging.getLogger(__name__)

# Defensive cap — a runaway deck shouldn't create thousands of slide rows.
_MAX_SLIDES = 300
_TITLE_MAX = 300


def parse_pptx(file_bytes: bytes) -> List[dict]:
    """Turn .pptx bytes into a list of {title, body_md, speaker_notes}.

    Slide order follows the deck's own order. Each slide's title comes from
    its title placeholder when present, otherwise the first line of text on
    the slide. Remaining text shapes become markdown bullets, and the slide's
    speaker notes (if any) are carried over verbatim.

    Raises ValueError on an unreadable file or a deck with no usable slides so
    the caller can return a clean 4xx.
    """
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as exc:  # noqa: BLE001 — any open failure is a bad upload
        logger.warning("python-pptx failed to open uploaded deck: %s", exc)
        raise ValueError(
            "Could not open the PowerPoint file. Make sure it's a valid "
            ".pptx (in Google Slides: File → Download → Microsoft PowerPoint)."
        ) from exc

    slides_out: List[dict] = []
    for idx, slide in enumerate(prs.slides):
        if idx >= _MAX_SLIDES:
            logger.info("Deck exceeded %s slides; truncating.", _MAX_SLIDES)
            break

        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:  # noqa: BLE001 — some layouts have no title placeholder
            title_shape = None

        title = ""
        if title_shape is not None and title_shape.has_text_frame:
            title = title_shape.text.strip()

        body_lines: List[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if title_shape is not None and shape.shape_id == title_shape.shape_id:
                continue
            for line in shape.text.splitlines():
                line = line.strip()
                if line:
                    body_lines.append(line)

        if not title:
            # No title placeholder — promote the first text line to the title.
            title = body_lines.pop(0) if body_lines else f"Slide {idx + 1}"

        body_md = "\n".join(f"- {line}" for line in body_lines)

        speaker_notes = ""
        try:
            if slide.has_notes_slide:
                speaker_notes = (
                    slide.notes_slide.notes_text_frame.text or ""
                ).strip()
        except Exception:  # noqa: BLE001 — malformed notes shouldn't fail the deck
            speaker_notes = ""

        slides_out.append({
            "title": title[:_TITLE_MAX],
            "body_md": body_md,
            "speaker_notes": speaker_notes,
        })

    if not slides_out:
        raise ValueError(
            "The presentation has no slides with readable text. If it's all "
            "images, the AI-generated option may suit it better."
        )
    return slides_out


def recommended_periods(total_slides: int, slides_per_period: int) -> int:
    """Ceiling division of slides into periods, never less than 1."""
    if total_slides <= 0 or slides_per_period <= 0:
        return 1
    return max(1, -(-total_slides // slides_per_period))
