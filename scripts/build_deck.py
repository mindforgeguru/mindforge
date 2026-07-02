#!/usr/bin/env python3
"""
Generate the MIND FORGE pitch deck (PowerPoint, ~12 slides).
Audience: senior engineering manager, security officer, CEO.

Run:  python3 scripts/build_deck.py
Out:  MindForge_Presentation.pptx  (repo root)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Brand logo (transparent MF mark) ──────────────────────────────────────────
_HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO = os.path.join(os.path.dirname(os.path.abspath(__file__)), "mindforge_logo.png")

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0B, 0x14, 0x26)   # background
NAVY_2    = RGBColor(0x12, 0x1F, 0x38)   # card
AMBER     = RGBColor(0xFF, 0x9F, 0x1C)   # forge accent
AMBER_DK  = RGBColor(0xE8, 0x8A, 0x0C)
WHITE     = RGBColor(0xF5, 0xF7, 0xFA)
GREY      = RGBColor(0x9A, 0xA8, 0xBD)
GREEN     = RGBColor(0x3D, 0xDC, 0x97)
TEAL      = RGBColor(0x4C, 0xC9, 0xF0)

# 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────
def bg(slide, color=NAVY):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def para(tf, text, size, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
         space_after=6, space_before=0, first=False, font="Calibri", italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return p


def kicker(slide, text, x=Inches(0.7), y=Inches(0.55)):
    bar = rect(slide, x, y + Inches(0.02), Inches(0.32), Inches(0.30), AMBER)
    tb, tf = textbox(slide, x + Inches(0.5), y, Inches(8), Inches(0.4))
    para(tf, text.upper(), 14, AMBER, bold=True, first=True)
    return bar


def heading(slide, text, x=Inches(0.7), y=Inches(0.95), w=Inches(12), size=34):
    tb, tf = textbox(slide, x, y, w, Inches(1.1))
    para(tf, text, size, WHITE, bold=True, first=True)
    return tb


def page_num(slide, n):
    tb, tf = textbox(slide, SW - Inches(1.1), SH - Inches(0.55), Inches(0.8), Inches(0.35))
    para(tf, f"{n:02d} / 12", 11, GREY, align=PP_ALIGN.RIGHT, first=True)


def logo(slide, x, y, size):
    """Place the transparent MF logo at (x, y) sized to `size` (square)."""
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, x, y, height=size, width=size)


def corner_logo(slide):
    """Small brand mark, top-right of content slides."""
    logo(slide, SW - Inches(1.15), Inches(0.5), Inches(0.6))


def chip(slide, x, y, text, color=TEAL):
    w = Inches(0.16 + 0.092 * len(text))
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.42))
    sp.adjustments[0] = 0.5
    sp.fill.solid()
    sp.fill.fore_color.rgb = NAVY_2
    sp.line.color.rgb = color
    sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = 0; tf.margin_bottom = 0
    para(tf, text, 12.5, color, bold=True, align=PP_ALIGN.CENTER, first=True)
    return x + w + Inches(0.18)


def card(slide, x, y, w, h, title, body, accent=AMBER, title_size=18, body_size=14):
    rect(slide, x, y, w, h, NAVY_2)
    rect(slide, x, y, Inches(0.09), h, accent)           # accent spine
    tb, tf = textbox(slide, x + Inches(0.32), y + Inches(0.26), w - Inches(0.55), h - Inches(0.5))
    para(tf, title, title_size, WHITE, bold=True, first=True, space_after=8)
    for line in body:
        para(tf, line, body_size, GREY, space_after=5)
    return tb


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
# accent side band
rect(s, 0, 0, Inches(0.22), SH, AMBER)
# brand logo
logo(s, Inches(0.85), Inches(1.05), Inches(1.25))

tb, tf = textbox(s, Inches(0.9), Inches(2.45), Inches(11), Inches(2.2))
para(tf, "MIND FORGE", 60, WHITE, bold=True, first=True, space_after=4)
para(tf, "An AI-powered school platform — tests from any PDF,", 22, GREY, space_after=2)
para(tf, "auto-graded, with real-time visibility for parents and admins.", 22, GREY)

tb, tf = textbox(s, Inches(0.9), Inches(5.7), Inches(11), Inches(1))
para(tf, "Teacher  ·  Student  ·  Parent  ·  Admin — one mobile app", 16, AMBER, bold=True, first=True, space_after=4)
para(tf, "Flutter + FastAPI  ·  Live / In Production  ·  May 2026", 13, GREY)

tb, tf = textbox(s, SW - Inches(4.2), SH - Inches(0.85), Inches(3.9), Inches(0.5))
para(tf, "Chinmay Jobanputra", 13, GREY, align=PP_ALIGN.RIGHT, first=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "The Problem")
heading(s, "Mid-sized schools run on a patchwork of disconnected tools")
y = Inches(2.15); h = Inches(1.95); w = Inches(2.85); gap = Inches(0.25); x = Inches(0.7)
cards = [
    ("Teachers", ["Hours hand-grading tests", "Paper attendance registers", "Chasing homework by hand"]),
    ("Students", ["Deadlines lost in WhatsApp", "Test prep disconnected", "from class material"]),
    ("Parents", ["Find out only at quarterly", "PTMs", "Paper receipts & report cards"]),
    ("Admins", ["Excel sheets for fees,", "timetables, year rollover", "No single source of truth"]),
]
cx = x
for t, b in cards:
    card(s, cx, y, w, h, t, b, accent=AMBER, title_size=18, body_size=13)
    cx += w + gap
tb, tf = textbox(s, Inches(0.7), Inches(4.6), Inches(12), Inches(1.2))
para(tf, "No single, affordable platform ties AI assessment, attendance, fees, and",
     19, WHITE, bold=True, first=True, space_after=2)
para(tf, "parent communication together for Indian secondary schools (grades 8–10).", 19, WHITE, bold=True)
page_num(s, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "The Solution")
heading(s, "One app. Four roles. One source of truth.")
tb, tf = textbox(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.7))
para(tf, "A teacher turns a PDF chapter into an auto-graded test in minutes — students attempt it on their",
     16, GREY, first=True, space_after=2)
para(tf, "phone, parents see results instantly, admins keep records, fees and timetables in one place.", 16, GREY)

feats = [
    ("AI Test Generation", "PDF → structured test in < 2 min", AMBER),
    ("Single-Attempt Tests", "Can't be gamed by app-killing", AMBER),
    ("Attendance & Timetable", "Per-period, holiday-aware", TEAL),
    ("Grades & Reports", "Auto-graded + PDF report cards", TEAL),
    ("Fees Module", "Structures, dues, payment history", GREEN),
    ("Broadcasts & Push", "FCM notices to school or grade", GREEN),
]
y = Inches(2.75); w = Inches(3.9); h = Inches(1.35); gx = Inches(0.25); gy = Inches(0.28)
x0 = Inches(0.7)
for i, (t, d, c) in enumerate(feats):
    col = i % 3; row = i // 3
    cx = x0 + col * (w + gx)
    cy = y + row * (h + gy)
    card(s, cx, cy, w, h, t, [d], accent=c, title_size=16, body_size=13)
page_num(s, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — FLAGSHIP DIFFERENTIATOR: AI TEST GENERATION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Flagship Differentiator")
heading(s, "From a PDF chapter to a live test in under 2 minutes")
# pipeline steps
steps = [
    ("1", "Upload", "Teacher uploads a PDF\nor image of a chapter"),
    ("2", "OCR", "PyMuPDF + Tesseract\nextract the text"),
    ("3", "Generate", "Gemini 2.5 Flash drafts\nMCQ + short-answer\n(Groq LLaMA 3.3 fallback)"),
    ("4", "Review", "Teacher edits &\npublishes"),
    ("5", "Auto-grade", "Objective graded instantly;\nsubjective on one screen"),
]
y = Inches(2.35); w = Inches(2.28); h = Inches(2.2); gap = Inches(0.18); x = Inches(0.7)
cx = x
for num, t, d in steps:
    rect(s, cx, y, w, h, NAVY_2)
    rect(s, cx, y, w, Inches(0.09), AMBER)
    cnum = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.28), y + Inches(0.3), Inches(0.55), Inches(0.55))
    cnum.fill.solid(); cnum.fill.fore_color.rgb = AMBER; cnum.line.fill.background(); cnum.shadow.inherit = False
    ntf = cnum.text_frame; ntf.margin_top = 0; ntf.margin_bottom = 0
    para(ntf, num, 20, NAVY, bold=True, align=PP_ALIGN.CENTER, first=True)
    tb, tf = textbox(s, cx + Inches(0.28), y + Inches(1.0), w - Inches(0.5), h - Inches(1.1))
    para(tf, t, 17, WHITE, bold=True, first=True, space_after=5)
    for line in d.split("\n"):
        para(tf, line, 12, GREY, space_after=2)
    if cx + w + gap < x + 5 * (w + gap) - gap:
        arr = textbox(s, cx + w - Inches(0.02), y + Inches(0.95), Inches(0.25), Inches(0.5))
        para(arr[1], "›", 26, AMBER, bold=True, align=PP_ALIGN.CENTER, first=True)
    cx += w + gap
tb, tf = textbox(s, Inches(0.7), Inches(4.95), Inches(12), Inches(1.2))
para(tf, "Why it matters:  no other Indian school information system does AI test generation this fast.",
     18, AMBER, bold=True, first=True, space_after=6)
para(tf, "Target: under 5 minutes from upload to publish. Test generation runs < 30s for a 5-page PDF.",
     15, GREY)
page_num(s, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — INTEGRITY: STRICT SINGLE-ATTEMPT
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Assessment Integrity")
heading(s, "Tests that can't be gamed")
left = Inches(0.7); colw = Inches(5.9)
card(s, left, Inches(2.15), colw, Inches(3.6), "Strict single-attempt enforcement", [
    "One attempt only — within a fixed 3-day window.",
    "",
    "Closing the app, switching screens, or running",
    "out of time auto-submits with the current state.",
    "",
    "A blank or abandoned attempt scores zero —",
    "there is no second try, no app-kill exploit.",
    "",
    "Submitted attempts are buffered if the network",
    "drops, so a flaky connection never loses work.",
], accent=AMBER, title_size=19, body_size=14)
rx = left + colw + Inches(0.4)
card(s, rx, Inches(2.15), colw, Inches(3.6), "Grading & reporting", [
    "Three grade types: online (auto-graded),",
    "offline (teacher-entered), and manual.",
    "",
    "Subjective answers graded from a single screen.",
    "",
    "Charts and percentage breakdowns on student",
    "and parent dashboards.",
    "",
    "Server-generated PDFs (ReportLab): test papers,",
    "answer keys, and result sheets.",
], accent=TEAL, title_size=19, body_size=14)
page_num(s, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE  (engineering manager)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "For the Engineering Lead")
heading(s, "Architecture: an async, horizontally-scalable stack")

def layer(slide, x, y, w, h, label, items, accent):
    rect(slide, x, y, w, h, NAVY_2)
    rect(slide, x, y, Inches(0.09), h, accent)
    tb, tf = textbox(slide, x + Inches(0.3), y + Inches(0.16), w - Inches(0.5), h - Inches(0.3))
    para(tf, label, 15, accent, bold=True, first=True, space_after=4)
    para(tf, items, 13, GREY)

x = Inches(0.7); w = Inches(11.95)
layer(s, x, Inches(2.0), w, Inches(0.95), "MOBILE  —  Flutter (iOS · Android · Web)",
      "Riverpod 2.5 state · go_router 13 · Dio 5.4 (auth + error interceptors) · Firebase Messaging · fl_chart · Material 3", TEAL)
layer(s, x, Inches(3.05), w, Inches(0.95), "EDGE  —  Nginx reverse proxy",
      "TLS termination · routing · static assets · multi-instance fronting", GREY)
layer(s, x, Inches(4.1), w, Inches(0.95), "API  —  FastAPI (async Python)",
      "SQLAlchemy 2.0 async + asyncpg · Alembic migrations · role-checked endpoints · WebSocket /ws/{user_id} for live updates", AMBER)
# data row — three boxes
dy = Inches(5.15); dw = Inches(3.83); dgap = Inches(0.23)
layer(s, x, dy, dw, Inches(1.0), "PostgreSQL", "Primary store · soft-deletes\n+ audit log", GREEN)
layer(s, x + dw + dgap, dy, dw, Inches(1.0), "Redis", "Token revocation +\nWebSocket pub/sub fanout", GREEN)
layer(s, x + 2*(dw + dgap), dy, dw, Inches(1.0), "MinIO (S3)", "4 buckets: tests, profiles,\nPDFs, DB files", GREEN)
tb, tf = textbox(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.6))
para(tf, "Async end-to-end + Redis pub/sub fanout → scales horizontally across instances.  Docker Compose for local & prod.",
     13.5, GREY, italic=True, first=True)
page_num(s, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SECURITY  (security officer)  ★
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "For the Security Officer")
heading(s, "Security: defence in depth")
items = [
    ("Authentication", "MPIN bcrypt-hashed (per-user salt). Designed for parents & younger students who struggle with passwords."),
    ("Token model", "Short-lived JWT (1 h access / 30 d refresh). Every token has a unique jti."),
    ("Revocation", "Redis-backed jti blocklist — logout instantly invalidates a token server-side."),
    ("Authorization", "Role-checked dependencies on every endpoint (admin / teacher / student / parent). Accounts need explicit admin approval before access."),
    ("Transport", "TLS everywhere; mobile app CA-pins Let's Encrypt to resist MITM."),
    ("Auditability", "Soft-deletes preserved + append-only audit log; Sentry tags errors to the acting user."),
]
y = Inches(2.0); w = Inches(5.85); h = Inches(1.55); gx = Inches(0.25); gy = Inches(0.18)
x0 = Inches(0.7)
for i, (t, d) in enumerate(items):
    col = i % 2; row = i // 2
    cx = x0 + col * (w + gx)
    cy = y + row * (h + gy)
    card(s, cx, cy, w, h, t, [d], accent=GREEN, title_size=16, body_size=13)
page_num(s, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PRIVACY & DATA PROTECTION  (security officer / CEO)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Privacy & Data Protection")
heading(s, "Built for minors' data — least privilege by default")
rows = [
    ("Data minimisation & masking", "PII (phone, email) is partially masked for non-privileged views — e.g. +91 98***3210, u***@example.com."),
    ("Strict tenancy", "Student data is isolated per school. A parent sees only their own linked child — read-only."),
    ("Least-privilege roles", "Each role sees exactly what its daily decisions need — nothing more. Enforced server-side, not just in the UI."),
    ("Right to erasure path", "Soft-delete + audit trail balances record-keeping with the ability to remove a user on request."),
    ("Published policy", "A written Privacy Policy ships with the app, covering collection, retention, and parental consent."),
]
y = Inches(2.05)
for t, d in rows:
    rect(s, Inches(0.7), y, Inches(11.95), Inches(0.86), NAVY_2)
    rect(s, Inches(0.7), y, Inches(0.09), Inches(0.86), TEAL)
    tb, tf = textbox(s, Inches(1.0), y + Inches(0.13), Inches(3.6), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, 15.5, WHITE, bold=True, first=True)
    tb, tf = textbox(s, Inches(4.75), y + Inches(0.13), Inches(7.6), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, d, 13.5, GREY, first=True)
    y += Inches(0.97)
page_num(s, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BUSINESS VALUE / MARKET  (CEO)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "For the CEO — Market & Value")
heading(s, "A wedge into a fragmented, underserved market")
# metric tiles
tiles = [
    ("8–10", "Grades targeted first;\nexpandable to 6–7, 11–12", AMBER),
    ("300–1,500", "Students per institution\n(mid-sized schools)", TEAL),
    ("< 5 min", "Upload → published test\n(was hours of grading)", GREEN),
    ("4-in-1", "Teacher, student, parent,\nadmin on one platform", AMBER),
]
y = Inches(2.05); w = Inches(2.85); h = Inches(1.85); gap = Inches(0.25); x = Inches(0.7)
cx = x
for big, sub, c in tiles:
    rect(s, cx, y, w, h, NAVY_2)
    rect(s, cx, y, w, Inches(0.09), c)
    tb, tf = textbox(s, cx + Inches(0.2), y + Inches(0.32), w - Inches(0.4), h - Inches(0.5))
    para(tf, big, 34, c, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=6)
    for line in sub.split("\n"):
        para(tf, line, 12.5, GREY, align=PP_ALIGN.CENTER, space_after=1)
    cx += w + gap
tb, tf = textbox(s, Inches(0.7), Inches(4.35), Inches(11.9), Inches(2.5))
para(tf, "The business case", 18, WHITE, bold=True, first=True, space_after=8)
para(tf, "•  Replaces paper registers, WhatsApp groups, and disconnected portals with one affordable subscription.",
     15, GREY, space_after=5)
para(tf, "•  AI assessment is the hook; attendance, fees, and parent visibility drive daily stickiness and renewals.",
     15, GREY, space_after=5)
para(tf, "•  Tracked metrics: teachers publishing ≥ 1 test/week, daily active parents, time-to-publish, YoY renewal, NPS > 40.",
     15, GREY, space_after=5)
para(tf, "•  Mobile-first and MPIN auth — built for the actual humans using it, not IT admins.",
     15, GREY)
page_num(s, 9)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — WHY IT WINS (DIFFERENTIATORS)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Why It Wins")
heading(s, "Five things competitors don't have together")
diffs = [
    ("AI test generation < 2 min", "From any PDF — no other Indian SIS does this."),
    ("Un-gameable assessments", "Strict single-attempt; app-killing won't help."),
    ("One source of truth", "Teacher, student, parent, admin share one record."),
    ("Mobile-first by design", "Built for the phone, not a retrofitted web SIS."),
    ("Auth for real users", "4-digit MPIN — not built for IT admins."),
]
y = Inches(2.1); w = Inches(11.95); h = Inches(0.82)
for i, (t, d) in enumerate(diffs):
    rect(s, Inches(0.7), y, w, h, NAVY_2)
    cnum = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), y + Inches(0.19), Inches(0.44), Inches(0.44))
    cnum.fill.solid(); cnum.fill.fore_color.rgb = AMBER; cnum.line.fill.background(); cnum.shadow.inherit = False
    ntf = cnum.text_frame; ntf.margin_top = 0; ntf.margin_bottom = 0
    para(ntf, str(i + 1), 17, NAVY, bold=True, align=PP_ALIGN.CENTER, first=True)
    tb, tf = textbox(s, Inches(1.7), y + Inches(0.11), Inches(4.5), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, 16, WHITE, bold=True, first=True)
    tb, tf = textbox(s, Inches(6.3), y + Inches(0.11), Inches(6.2), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, d, 13.5, GREY, first=True)
    y += h + Inches(0.13)
page_num(s, 10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — STATUS & ROADMAP
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Status & Roadmap")
heading(s, "Live today — with a clear next horizon")
# Done column
card(s, Inches(0.7), Inches(2.1), Inches(5.9), Inches(4.1), "Shipped & in production", [
    "✓  All four role dashboards live",
    "✓  AI test generation + auto-grading",
    "✓  Attendance, timetable, holiday-aware",
    "✓  Fees module + parent visibility",
    "✓  Broadcasts + FCM push notifications",
    "✓  Security hardening (CA pinning, JWT revocation)",
    "✓  April 2026 performance pass — complete",
    "✓  Android release signing — done",
], accent=GREEN, title_size=19, body_size=14.5)
card(s, Inches(6.8), Inches(2.1), Inches(5.85), Inches(4.1), "Next", [
    "•  iOS App Store submission (Apple Dev enrolment)",
    "•  Offline-first test attempts with sync",
    "•  Parent–teacher 1:1 chat",
    "•  Multi-language (Hindi + regional)",
    "•  Admin analytics: cohort & trend reporting",
    "•  Subject expansion: grades 6–7 and 11–12",
    "•  Gamification: student streaks & badges",
], accent=AMBER, title_size=19, body_size=14.5)
page_num(s, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, Inches(0.22), SH, AMBER)
logo(s, Inches(0.9), Inches(0.7), Inches(1.15))
tb, tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.5))
para(tf, "MIND FORGE", 46, WHITE, bold=True, first=True, space_after=10)
para(tf, "AI assessment, attendance, fees and parent visibility —", 22, GREY, space_after=2)
para(tf, "in one mobile app, secured for minors' data.", 22, GREY)
tb, tf = textbox(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(1.6))
para(tf, "Built end-to-end:  Flutter + FastAPI · Gemini AI · PostgreSQL / Redis / MinIO · Docker.",
     16, AMBER, bold=True, first=True, space_after=14)
para(tf, "Thank you — happy to take questions, or run a live demo.", 18, WHITE, bold=True)
tb, tf = textbox(s, Inches(0.9), SH - Inches(0.85), Inches(11.5), Inches(0.5))
para(tf, "Chinmay Jobanputra  ·  chinmay1975@gmail.com", 13, GREY, first=True)


# ── Brand mark on content slides (2–11) ───────────────────────────────────────
for sl in list(prs.slides)[1:11]:
    corner_logo(sl)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "MindForge_Presentation.pptx"
prs.save(out)
print(f"Saved {out} — {len(prs.slides._sldIdLst)} slides")
